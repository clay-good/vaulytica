"""Gmail attachment scanner for PII detection."""

import base64
import io
import time
from dataclasses import dataclass, field
from datetime import datetime, timezone, timedelta
from typing import List, Optional, Dict, Any, Iterator

import structlog
from googleapiclient.errors import HttpError

from vaulytica.core.auth.client import GoogleWorkspaceClient
from vaulytica.core.detectors.pii_detector import PIIDetector, PIIDetectionResult

# Document parsing imports
try:
    from pypdf import PdfReader
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from openpyxl import load_workbook
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

logger = structlog.get_logger(__name__)


@dataclass
class EmailAttachment:
    """Represents an email attachment."""

    message_id: str
    attachment_id: str
    filename: str
    mime_type: str
    size: int
    user_email: str
    date: datetime
    subject: str
    sender: str
    recipients: List[str] = field(default_factory=list)  # To, Cc, Bcc
    external_recipients: List[str] = field(default_factory=list)  # Non-domain recipients
    is_sent_externally: bool = False


@dataclass
class AttachmentScanResult:
    """Result of scanning an attachment."""

    attachment: EmailAttachment
    pii_result: Optional[PIIDetectionResult] = None
    risk_score: int = 0
    scanned: bool = False
    error: Optional[str] = None


@dataclass
class GmailScanResult:
    """Results from Gmail scanning."""

    total_messages: int = 0
    total_attachments: int = 0
    attachments_scanned: int = 0
    attachments_with_pii: int = 0
    results: List[AttachmentScanResult] = field(default_factory=list)


class GmailScannerError(Exception):
    """Raised when Gmail scanning fails."""

    pass


class GmailScanner:
    """Scans Gmail attachments for PII."""

    # Scannable MIME types
    SCANNABLE_TYPES = [
        "text/plain",
        "text/html",
        "text/csv",
        "application/pdf",
        "application/msword",  # .doc (not fully supported)
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # .docx
        "application/vnd.ms-excel",  # .xls (not fully supported)
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # .xlsx
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
    ]

    def __init__(
        self,
        client: GoogleWorkspaceClient,
        domain: str,
    ):
        """Initialize Gmail scanner.

        Args:
            client: GoogleWorkspaceClient instance
            domain: Organization domain
        """
        self.client = client
        self.domain = domain
        self.pii_detector = PIIDetector()

        logger.info("gmail_scanner_initialized", domain=domain)

    def scan_user_attachments(
        self,
        user_email: str,
        days_back: int = 30,
        max_messages: int = 100,
        external_only: bool = False,
        max_attachments: Optional[int] = None,
    ) -> GmailScanResult:
        """Scan attachments for a user with enhanced performance.

        Args:
            user_email: User email to scan
            days_back: Number of days to look back
            max_messages: Maximum messages to scan
            external_only: Only scan emails sent to external recipients
            max_attachments: Maximum number of attachments to scan (for performance testing)

        Returns:
            GmailScanResult

        Raises:
            ValueError: If invalid parameters provided
        """
        # Input validation
        if not isinstance(user_email, str) or not user_email:
            raise ValueError("user_email must be a non-empty string")
        if days_back < 1:
            raise ValueError("days_back must be at least 1")
        if max_messages < 1:
            raise ValueError("max_messages must be at least 1")
        if max_attachments is not None and (not isinstance(max_attachments, int) or max_attachments < 1):
            raise ValueError("max_attachments must be a positive integer")

        logger.info(
            "scanning_user_attachments",
            user_email=user_email,
            days_back=days_back,
            max_messages=max_messages,
            external_only=external_only,
            max_attachments=max_attachments,
        )
        scan_start_time = time.time()

        result = GmailScanResult()
        failed_messages = []
        attachment_count = 0

        # Build query for messages with attachments
        after_date = datetime.now(timezone.utc) - timedelta(days=days_back)
        query = f"has:attachment after:{after_date.strftime('%Y/%m/%d')}"

        # Get messages
        messages = list(
            self._list_messages(user_email, query, max_results=max_messages)
        )

        result.total_messages = len(messages)

        # Scan each message
        for message_id in messages:
            try:
                attachments = self._get_message_attachments(user_email, message_id)
                result.total_attachments += len(attachments)

                for attachment in attachments:
                    # Filter by external_only if requested
                    if external_only and not attachment.is_sent_externally:
                        continue

                    scan_result = self._scan_attachment(user_email, attachment)
                    result.results.append(scan_result)

                    if scan_result.scanned:
                        result.attachments_scanned += 1

                    if scan_result.pii_result and scan_result.pii_result.has_pii:
                        result.attachments_with_pii += 1

                    attachment_count += 1

                    # Log progress every 50 attachments
                    if attachment_count % 50 == 0:
                        logger.info(
                            "gmail_scan_progress",
                            scanned=attachment_count,
                            with_pii=result.attachments_with_pii,
                            messages_processed=len(failed_messages) + 1,
                        )

                    # Check max_attachments limit
                    if max_attachments and attachment_count >= max_attachments:
                        logger.info("max_attachments_limit_reached", max_attachments=max_attachments)
                        break

                # Check if we hit max_attachments limit
                if max_attachments and attachment_count >= max_attachments:
                    break

            except HttpError as e:
                if e.resp.status == 403:
                    logger.error(
                        "insufficient_permissions_to_scan_message",
                        user_email=user_email,
                        message_id=message_id,
                        error=str(e),
                    )
                else:
                    logger.warning(
                        "failed_to_scan_message",
                        user_email=user_email,
                        message_id=message_id,
                        error=str(e),
                    )
                failed_messages.append({
                    "message_id": message_id,
                    "error": str(e)
                })
            except Exception as e:
                logger.warning(
                    "failed_to_scan_message",
                    user_email=user_email,
                    message_id=message_id,
                    error=str(e),
                )
                failed_messages.append({
                    "message_id": message_id,
                    "error": str(e)
                })

        # Calculate scan duration
        scan_duration = time.time() - scan_start_time

        logger.info(
            "gmail_scan_complete",
            user_email=user_email,
            total_messages=result.total_messages,
            attachments_scanned=result.attachments_scanned,
            attachments_with_pii=result.attachments_with_pii,
            failed_messages=len(failed_messages),
            scan_duration_seconds=round(scan_duration, 2),
        )

        # Log warning if many messages failed
        if failed_messages and len(failed_messages) > 5:
            logger.warning(
                "many_messages_failed_scan",
                failed_count=len(failed_messages),
                sample_errors=failed_messages[:3]
            )

        return result

    def _list_messages(
        self,
        user_email: str,
        query: str,
        max_results: int = 100,
    ) -> Iterator[str]:
        """List messages matching query.

        Args:
            user_email: User email
            query: Gmail search query
            max_results: Maximum results

        Yields:
            Message IDs
        """
        page_token = None
        count = 0

        try:
            while count < max_results:
                try:
                    response = (
                        self.client.gmail.users()
                        .messages()
                        .list(
                            userId=user_email,
                            q=query,
                            maxResults=min(100, max_results - count),
                            pageToken=page_token,
                        )
                        .execute()
                    )

                    messages = response.get("messages", [])

                    for message in messages:
                        yield message["id"]
                        count += 1

                    page_token = response.get("nextPageToken")
                    if not page_token:
                        break

                    time.sleep(0.1)  # Rate limiting

                except HttpError as e:
                    if e.resp.status == 429:  # Rate limit
                        logger.warning("rate_limit_hit_retrying")
                        time.sleep(5)
                        continue
                    else:
                        raise

        except Exception as e:
            logger.error("failed_to_list_messages", user_email=user_email, error=str(e))
            raise GmailScannerError(f"Failed to list messages: {e}")

    def _get_message_attachments(
        self,
        user_email: str,
        message_id: str,
    ) -> List[EmailAttachment]:
        """Get attachments from a message.

        Args:
            user_email: User email
            message_id: Message ID

        Returns:
            List of EmailAttachment objects
        """
        try:
            message = (
                self.client.gmail.users()
                .messages()
                .get(userId=user_email, id=message_id, format="full")
                .execute()
            )

            attachments = []

            # Get message metadata
            headers = message.get("payload", {}).get("headers", [])
            subject = next(
                (h["value"] for h in headers if h["name"].lower() == "subject"),
                "No Subject",
            )
            sender = next(
                (h["value"] for h in headers if h["name"].lower() == "from"),
                "Unknown",
            )
            date_str = next(
                (h["value"] for h in headers if h["name"].lower() == "date"),
                None,
            )

            # Parse date
            if date_str:
                from email.utils import parsedate_to_datetime

                date = parsedate_to_datetime(date_str)
            else:
                date = datetime.now(timezone.utc)

            # Extract recipients (To, Cc, Bcc)
            recipients = []
            for header in headers:
                if header["name"].lower() in ["to", "cc", "bcc"]:
                    # Parse email addresses from header value
                    # Format: "Name <email@domain.com>, Name2 <email2@domain.com>"
                    import re
                    email_pattern = r'[\w\.-]+@[\w\.-]+'
                    emails = re.findall(email_pattern, header["value"])
                    recipients.extend(emails)

            # Determine external recipients
            external_recipients = []
            for recipient in recipients:
                recipient_domain = recipient.split("@")[-1].lower()
                if recipient_domain != self.domain.lower():
                    external_recipients.append(recipient)

            is_sent_externally = len(external_recipients) > 0

            # Extract attachments
            parts = self._get_message_parts(message.get("payload", {}))

            for part in parts:
                if part.get("filename") and part.get("body", {}).get("attachmentId"):
                    attachment = EmailAttachment(
                        message_id=message_id,
                        attachment_id=part["body"]["attachmentId"],
                        filename=part["filename"],
                        mime_type=part.get("mimeType", ""),
                        size=part["body"].get("size", 0),
                        user_email=user_email,
                        date=date,
                        subject=subject,
                        sender=sender,
                        recipients=recipients,
                        external_recipients=external_recipients,
                        is_sent_externally=is_sent_externally,
                    )
                    attachments.append(attachment)

            return attachments

        except HttpError as e:
            logger.error(
                "failed_to_get_message",
                user_email=user_email,
                message_id=message_id,
                error=str(e),
            )
            return []

    def _get_message_parts(self, payload: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Recursively get all message parts.

        Args:
            payload: Message payload

        Returns:
            List of parts
        """
        parts = []

        if "parts" in payload:
            for part in payload["parts"]:
                parts.extend(self._get_message_parts(part))
        else:
            parts.append(payload)

        return parts

    def _scan_attachment(
        self,
        user_email: str,
        attachment: EmailAttachment,
    ) -> AttachmentScanResult:
        """Scan an attachment for PII.

        Args:
            user_email: User email
            attachment: EmailAttachment object

        Returns:
            AttachmentScanResult
        """
        result = AttachmentScanResult(attachment=attachment)

        # Check if scannable
        if attachment.mime_type not in self.SCANNABLE_TYPES:
            result.scanned = False
            return result

        # Check size (skip large files)
        if attachment.size > 10 * 1024 * 1024:  # 10MB
            result.scanned = False
            result.error = "File too large"
            return result

        try:
            # Get attachment data
            attachment_data = (
                self.client.gmail.users()
                .messages()
                .attachments()
                .get(
                    userId=user_email,
                    messageId=attachment.message_id,
                    id=attachment.attachment_id,
                )
                .execute()
            )

            # Decode data
            data = base64.urlsafe_b64decode(attachment_data["data"])

            # Extract text based on MIME type
            text = self._extract_text(data, attachment.mime_type)

            if text:
                # Scan for PII
                pii_result = self.pii_detector.detect(text)
                result.pii_result = pii_result
                result.risk_score = self._calculate_risk_score(pii_result)
                result.scanned = True
            else:
                result.scanned = False
                result.error = "Could not extract text"

        except Exception as e:
            logger.error(
                "failed_to_scan_attachment",
                user_email=user_email,
                attachment_id=attachment.attachment_id,
                error=str(e),
            )
            result.scanned = False
            result.error = str(e)

        return result

    def _extract_text(self, data: bytes, mime_type: str) -> Optional[str]:
        """Extract text from attachment data.

        Args:
            data: Attachment data
            mime_type: MIME type

        Returns:
            Extracted text or None
        """
        try:
            # Text files
            if mime_type.startswith("text/"):
                return data.decode("utf-8", errors="ignore")

            # PDF files
            elif mime_type == "application/pdf":
                return self._extract_pdf_text(data)

            # Word documents (.docx)
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                return self._extract_docx_text(data)

            # Old Word documents (.doc)
            elif mime_type == "application/msword":
                logger.warning("old_word_format_not_supported", mime_type=mime_type)
                return None

            # Excel spreadsheets (.xlsx)
            elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                return self._extract_xlsx_text(data)

            # Old Excel (.xls)
            elif mime_type == "application/vnd.ms-excel":
                logger.warning("old_excel_format_not_supported", mime_type=mime_type)
                return None

            # PowerPoint (.pptx)
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return self._extract_pptx_text(data)

            else:
                logger.debug("unsupported_mime_type", mime_type=mime_type)
                return None

        except Exception as e:
            logger.error("failed_to_extract_text", mime_type=mime_type, error=str(e))
            return None

    def _extract_pdf_text(self, data: bytes) -> Optional[str]:
        """Extract text from PDF.

        Args:
            data: PDF file data

        Returns:
            Extracted text or None
        """
        if not HAS_PDF:
            logger.warning("pypdf2_not_installed")
            return None

        try:
            pdf_file = io.BytesIO(data)
            reader = PdfReader(pdf_file)

            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)

            return "\n".join(text_parts) if text_parts else None

        except Exception as e:
            logger.error("failed_to_extract_pdf", error=str(e))
            return None

    def _extract_docx_text(self, data: bytes) -> Optional[str]:
        """Extract text from Word document (.docx).

        Args:
            data: DOCX file data

        Returns:
            Extracted text or None
        """
        if not HAS_DOCX:
            logger.warning("python_docx_not_installed")
            return None

        try:
            docx_file = io.BytesIO(data)
            doc = Document(docx_file)

            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    text_parts.append(paragraph.text)

            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text_parts.append(cell.text)

            return "\n".join(text_parts) if text_parts else None

        except Exception as e:
            logger.error("failed_to_extract_docx", error=str(e))
            return None

    def _extract_xlsx_text(self, data: bytes) -> Optional[str]:
        """Extract text from Excel spreadsheet (.xlsx).

        Args:
            data: XLSX file data

        Returns:
            Extracted text or None
        """
        if not HAS_XLSX:
            logger.warning("openpyxl_not_installed")
            return None

        try:
            xlsx_file = io.BytesIO(data)
            workbook = load_workbook(xlsx_file, data_only=True)

            text_parts = []
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell_value in row:
                        if cell_value is not None:
                            text_parts.append(str(cell_value))

            return "\n".join(text_parts) if text_parts else None

        except Exception as e:
            logger.error("failed_to_extract_xlsx", error=str(e))
            return None

    def _extract_pptx_text(self, data: bytes) -> Optional[str]:
        """Extract text from PowerPoint presentation (.pptx).

        Args:
            data: PPTX file data

        Returns:
            Extracted text or None
        """
        if not HAS_PPTX:
            logger.warning("python_pptx_not_installed")
            return None

        try:
            pptx_file = io.BytesIO(data)
            presentation = Presentation(pptx_file)

            text_parts = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)

            return "\n".join(text_parts) if text_parts else None

        except Exception as e:
            logger.error("failed_to_extract_pptx", error=str(e))
            return None

    def _calculate_risk_score(self, pii_result: PIIDetectionResult) -> int:
        """Calculate risk score for attachment.

        Args:
            pii_result: PII detection result

        Returns:
            Risk score (0-100)
        """
        if not pii_result.has_pii:
            return 0

        score = 0

        # Base score for having PII
        score += 30

        # Add points for each PII type
        score += len(pii_result.findings) * 10

        # Add points for high-confidence findings
        high_confidence = [f for f in pii_result.findings if f.confidence >= 0.8]
        score += len(high_confidence) * 5

        return min(100, score)

